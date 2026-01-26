#!/usr/bin/env python3
"""
Script pour générer un slide PowerPoint à partir de la synthèse du projet.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx n'est pas installé.")
    print("Installez-le avec: pip install python-pptx")
    exit(1)

# Créer une présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Ajouter un slide avec layout vide
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide

# Titre principal
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "🚀 OVH Customer Feedbacks Tracker"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(102, 126, 234)  # #667eea
title_para.alignment = PP_ALIGN.CENTER

# Sous-titre
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Synthèse ROI - Développement avec VibeCoding"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = RGBColor(118, 75, 162)  # #764ba2
subtitle_para.alignment = PP_ALIGN.CENTER

# Box highlight - Projet développé par un Ingénieur Généraliste
highlight_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.2))
highlight_frame = highlight_box.text_frame
highlight_frame.text = "Projet développé par un Ingénieur Généraliste\n\n52 388 lignes\nen 3-4 semaines"
highlight_para = highlight_frame.paragraphs[0]
highlight_para.font.size = Pt(20)
highlight_para.font.bold = True
highlight_para.font.color.rgb = RGBColor(255, 255, 255)
highlight_para.alignment = PP_ALIGN.CENTER
highlight_shape = highlight_box
highlight_shape.fill.solid()
highlight_shape.fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea
highlight_shape.line.color.rgb = RGBColor(102, 126, 234)

# Stats grid - 3 colonnes
stats_y = 3.2
stats_width = 2.5
stats_height = 1.2
stats_gap = 0.3

# Stat 1: Gain vs Junior
stat1_box = slide.shapes.add_textbox(Inches(0.8), Inches(stats_y), Inches(stats_width), Inches(stats_height))
stat1_frame = stat1_box.text_frame
stat1_frame.text = "90-93%\nGain vs Dev Junior"
stat1_para = stat1_frame.paragraphs[0]
stat1_para.font.size = Pt(32)
stat1_para.font.bold = True
stat1_para.font.color.rgb = RGBColor(102, 126, 234)
stat1_para.alignment = PP_ALIGN.CENTER

# Stat 2: Gain vs Expérimenté
stat2_box = slide.shapes.add_textbox(Inches(0.8 + stats_width + stats_gap), Inches(stats_y), Inches(stats_width), Inches(stats_height))
stat2_frame = stat2_box.text_frame
stat2_frame.text = "75-85%\nGain vs Dev Expérimenté"
stat2_para = stat2_frame.paragraphs[0]
stat2_para.font.size = Pt(32)
stat2_para.font.bold = True
stat2_para.font.color.rgb = RGBColor(102, 126, 234)
stat2_para.alignment = PP_ALIGN.CENTER

# Stat 3: Productivité
stat3_box = slide.shapes.add_textbox(Inches(0.8 + 2 * (stats_width + stats_gap)), Inches(stats_y), Inches(stats_width), Inches(stats_height))
stat3_frame = stat3_box.text_frame
stat3_frame.text = "10-14x\nProductivité"
stat3_para = stat3_frame.paragraphs[0]
stat3_para.font.size = Pt(32)
stat3_para.font.bold = True
stat3_para.font.color.rgb = RGBColor(102, 126, 234)
stat3_para.alignment = PP_ALIGN.CENTER

# Tableau comparatif
table_y = 4.6
table_width = 9
table_height = 1.8

# Créer un tableau avec 4 colonnes et 4 lignes
table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(table_y), Inches(table_width), Inches(table_height)).table

# En-têtes
headers = ["Profil", "Durée", "Coût", "Productivité"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(102, 126, 234)
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.alignment = PP_ALIGN.CENTER

# Ligne 1: Dev Junior
row1_data = ["👨‍💻 Dev Junior (sans IA)", "6-9 mois", "18k-40.5k€", "~150-200 lignes/jour"]
for i, data in enumerate(row1_data):
    cell = table.cell(1, i)
    cell.text = data
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER

# Ligne 2: Dev Expérimenté
row2_data = ["👨‍💼 Dev Expérimenté (sans IA)", "4-6 mois", "20k-42k€", "~250 lignes/jour"]
for i, data in enumerate(row2_data):
    cell = table.cell(2, i)
    cell.text = data
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER

# Ligne 3: Ingénieur Généraliste (mise en évidence)
row3_data = ["🔧 Ingénieur Généraliste (avec VibeCoding)", "3-4 semaines", "20-30€ + 1 mois", "~2 500-3 500 lignes/jour"]
for i, data in enumerate(row3_data):
    cell = table.cell(3, i)
    cell.text = data
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(227, 242, 253)  # Bleu clair
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER

# ROI Box
roi_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
roi_frame = roi_box.text_frame
roi_frame.text = "💰 ROI Global: ~500-2 000x | Économie: 35 000-80 000€ | Investissement: 60-90€"
roi_para = roi_frame.paragraphs[0]
roi_para.font.size = Pt(18)
roi_para.font.bold = True
roi_para.font.color.rgb = RGBColor(46, 125, 50)  # Vert foncé
roi_para.alignment = PP_ALIGN.CENTER
roi_shape = roi_box
roi_shape.fill.solid()
roi_shape.fill.fore_color.rgb = RGBColor(232, 245, 233)  # Vert clair
roi_shape.line.color.rgb = RGBColor(76, 175, 80)

# Footer
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.3))
footer_frame = footer_box.text_frame
footer_frame.text = "Projet développé 100% avec VibeCoding (Cursor AI) | Janvier 2026 | Version 1.0.8"
footer_para = footer_frame.paragraphs[0]
footer_para.font.size = Pt(10)
footer_para.font.color.rgb = RGBColor(128, 128, 128)
footer_para.alignment = PP_ALIGN.CENTER

# Sauvegarder
output_file = "docs/SYNTHESE_ROI_PROJET.pptx"
prs.save(output_file)
print(f"Slide PowerPoint cree: {output_file}")

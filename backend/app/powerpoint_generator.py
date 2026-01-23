"""PowerPoint report generator for OVH feedback analytics."""
import io
import logging
import re
from typing import List, Dict, Optional
from datetime import datetime
from collections import Counter, defaultdict
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from PIL import Image
    PPTX_AVAILABLE = True
    MISSING_DEPENDENCIES = []
except ImportError as e:
    PPTX_AVAILABLE = False
    # Determine which dependency is missing
    missing = []
    try:
        import pptx
    except ImportError:
        missing.append("python-pptx")
    try:
        import matplotlib
    except ImportError:
        missing.append("matplotlib")
    try:
        from PIL import Image
    except ImportError:
        missing.append("Pillow")
    MISSING_DEPENDENCIES = missing if missing else ["python-pptx", "matplotlib", "Pillow"]

logger = logging.getLogger(__name__)


def generate_powerpoint_report(
    posts: List[Dict],
    filters: Dict,
    recommended_actions: List[Dict],
    stats: Dict,
    llm_analysis: Optional[str] = None,
    chart_images: Optional[Dict[str, bytes]] = None,
    improvements_analysis: Optional[List[Dict]] = None
) -> bytes:
    """
    Generate a PowerPoint report with charts and insights.
    
    Args:
        posts: List of post dictionaries
        filters: Current filter settings
        recommended_actions: List of recommended actions
        stats: Statistics dictionary
        llm_analysis: Optional LLM-generated analysis text
    
    Returns:
        bytes: PowerPoint file as bytes
    """
    if not PPTX_AVAILABLE:
        deps = ", ".join(MISSING_DEPENDENCIES)
        raise RuntimeError(
            f"PowerPoint generation requires {deps}. "
            f"Install with: pip install {' '.join(MISSING_DEPENDENCIES)} "
            f"or install all dependencies: pip install -r requirements.txt"
        )
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Build filter description
    filter_parts = []
    if filters.get('dateFrom') or filters.get('dateTo'):
        date_range = f"{filters.get('dateFrom', 'Start')} to {filters.get('dateTo', 'End')}"
        filter_parts.append(f"Period: {date_range}")
    if filters.get('search'):
        filter_parts.append(f"Search: \"{filters['search']}\"")
    if filters.get('sentiment') and filters['sentiment'] != 'all':
        filter_parts.append(f"Sentiment: {filters['sentiment']}")
    if filters.get('language') and filters['language'] != 'all':
        filter_parts.append(f"Language: {filters['language']}")
    if filters.get('source') and filters['source'] != 'all':
        filter_parts.append(f"Source: {filters['source']}")
    if filters.get('product') and filters['product'] != 'all':
        filter_parts.append(f"Product: {filters['product']}")
    
    filter_text = " | ".join(filter_parts) if filter_parts else "All posts"
    
    # Adjust title based on report type
    if improvements_analysis:
        title.text = "OVH Improvements Opportunities Report"
    else:
        title.text = "OVH Customer Feedback Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}\nBased on {len(posts)} feedback posts\n{filter_text}"
    
    # Slide 2: AI Analysis (Page 1) - Beautifully formatted
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "AI-Powered Analysis"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 153, 255)  # OVH blue
    
    # Add filter info box
    from pptx.enum.shapes import MSO_SHAPE
    filter_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
    filter_box.fill.solid()
    filter_box.fill.fore_color.rgb = RGBColor(245, 247, 250)  # Light gray
    filter_box.line.color.rgb = RGBColor(200, 200, 200)
    filter_box.line.width = Pt(1)
    
    filter_textbox = slide.shapes.add_textbox(Inches(0.6), Inches(1), Inches(8.8), Inches(0.3))
    filter_frame = filter_textbox.text_frame
    filter_frame.text = f"📊 Analysis Period: {filter_text}"
    filter_para = filter_frame.paragraphs[0]
    filter_para.font.size = Pt(12)
    filter_para.font.color.rgb = RGBColor(100, 100, 100)
    
    # Add stats summary
    stats_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    stats_box.line.color.rgb = RGBColor(0, 212, 255)
    stats_box.line.width = Pt(2)
    
    stats_textbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(0.6))
    stats_frame = stats_textbox.text_frame
    stats_frame.text = f"📈 Total Posts: {stats.get('total', len(posts))} | "
    stats_para = stats_frame.paragraphs[0]
    stats_para.font.size = Pt(14)
    stats_para.font.bold = True
    
    p = stats_frame.add_paragraph()
    p.text = f"✅ Positive: {stats.get('positive', 0)} | "
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(34, 197, 94)  # Green
    
    p = stats_frame.add_paragraph()
    p.text = f"❌ Negative: {stats.get('negative', 0)} | "
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(239, 68, 68)  # Red
    
    p = stats_frame.add_paragraph()
    p.text = f"⚪ Neutral: {stats.get('neutral', 0)}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(156, 163, 175)  # Gray
    
    # Format stats as inline
    stats_frame.paragraphs[0].text = f"📈 Total: {stats.get('total', len(posts))} | ✅ Positive: {stats.get('positive', 0)} | ❌ Negative: {stats.get('negative', 0)} | ⚪ Neutral: {stats.get('neutral', 0)}"
    stats_frame.paragraphs[0].font.size = Pt(14)
    stats_frame.paragraphs[0].font.bold = True
    # Remove extra paragraphs
    while len(stats_frame.paragraphs) > 1:
        p = stats_frame.paragraphs[1]
        stats_frame._element.remove(p._element)
    
    # Add LLM Analysis in a beautiful box
    analysis_top = Inches(2.5)
    analysis_height = Inches(4.0)  # Reduced to leave room for actions if needed
    
    if llm_analysis:
        analysis_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), analysis_top, Inches(9), analysis_height)
        analysis_box.fill.solid()
        analysis_box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
        analysis_box.line.color.rgb = RGBColor(0, 153, 255)  # OVH blue border
        analysis_box.line.width = Pt(3)
        
        analysis_textbox = slide.shapes.add_textbox(Inches(0.7), analysis_top + Inches(0.3), Inches(8.6), analysis_height - Inches(0.6))
        analysis_frame = analysis_textbox.text_frame
        analysis_frame.word_wrap = True
        
        # Parse and format LLM analysis
        analysis_lines = llm_analysis.split('\n')
        bullet_points = []
        for line in analysis_lines:
            line = line.strip()
            if line:
                line = line.lstrip('•-*').strip()
                if line:
                    bullet_points.append(line)
        
        # Add title
        analysis_frame.text = "💡 Key Insights"
        title_para = analysis_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 153, 255)
        title_para.space_after = Pt(12)
        
        # Add insights as formatted bullets
        # If improvements_analysis is provided, use it instead of llm_analysis
        if improvements_analysis and len(improvements_analysis) > 0:
            # Use improvements insights
            for i, insight in enumerate(improvements_analysis[:5]):  # Max 5 insights
                p = analysis_frame.add_paragraph()
                title_text = insight.get('title', '') if isinstance(insight, dict) else getattr(insight, 'title', '')
                desc_text = insight.get('description', '') if isinstance(insight, dict) else getattr(insight, 'description', '')
                if title_text and desc_text:
                    p.text = f"• {title_text}: {desc_text[:100]}..." if len(desc_text) > 100 else f"• {title_text}: {desc_text}"
                elif title_text:
                    p.text = f"• {title_text}"
                elif desc_text:
                    p.text = f"• {desc_text[:120]}..." if len(desc_text) > 120 else f"• {desc_text}"
                else:
                    continue
                p.level = 0
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_after = Pt(8)
                p.space_before = Pt(4)
        else:
            # Use LLM analysis bullets
            for i, point in enumerate(bullet_points[:5]):  # Max 5 insights
                p = analysis_frame.add_paragraph()
                p.text = f"• {point}"
                p.level = 0
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_after = Pt(8)
                p.space_before = Pt(4)
    
    # Add recommended actions if available (only if there's space)
    if recommended_actions and llm_analysis:
        # Adjust analysis height to leave room for actions
        analysis_height = Inches(3.5)
        actions_top = analysis_top + analysis_height + Inches(0.2)
        actions_height = Inches(1.0)
        
        # Check if actions fit on slide
        if actions_top + actions_height <= Inches(7.5):
            actions_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), actions_top, Inches(9), actions_height)
            actions_box.fill.solid()
            actions_box.fill.fore_color.rgb = RGBColor(255, 250, 240)  # Light beige
            actions_box.line.color.rgb = RGBColor(245, 158, 11)  # Orange border
            actions_box.line.width = Pt(2)
            
            actions_textbox = slide.shapes.add_textbox(Inches(0.7), actions_top + Inches(0.1), Inches(8.6), actions_height - Inches(0.2))
            actions_frame = actions_textbox.text_frame
            actions_frame.word_wrap = True
            actions_frame.text = "🎯 Recommended Actions"
            
            title_para = actions_frame.paragraphs[0]
            title_para.font.size = Pt(14)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(245, 158, 11)
            title_para.space_after = Pt(4)
            
            for action in recommended_actions[:2]:  # Max 2 actions to fit
                p = actions_frame.add_paragraph()
                p.text = f"• {action.get('text', 'N/A')[:80]}..." if len(action.get('text', '')) > 80 else f"• {action.get('text', 'N/A')}"
                p.level = 0
                p.font.size = Pt(10)
                priority = action.get('priority', 'medium').lower()
                if priority == 'high':
                    p.font.color.rgb = RGBColor(239, 68, 68)
                elif priority == 'medium':
                    p.font.color.rgb = RGBColor(245, 158, 11)
                else:
                    p.font.color.rgb = RGBColor(34, 197, 94)
                p.space_after = Pt(2)
    
    # Slide 3: Charts (Page 2) - Clean layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Data Visualization"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 153, 255)  # OVH blue
    
    # Add charts in a clean grid layout
    chart_images_dict = chart_images or {}
    
    # Calculate positions for 2x2 grid (or 3 charts)
    chart_width = Inches(4.2)
    chart_height = Inches(2.8)
    chart_spacing = Inches(0.3)
    
    chart_positions = [
        ('timeline', Inches(0.5), Inches(1.2), chart_width, chart_height),
        ('source', Inches(5.1), Inches(1.2), chart_width, chart_height),
        ('sentiment', Inches(2.8), Inches(4.3), chart_width, chart_height)
    ]
    
    charts_added = 0
    for chart_type, left, top, width, height in chart_positions:
        if chart_type in chart_images_dict and chart_images_dict[chart_type]:
            try:
                # Add chart
                chart_pic = slide.shapes.add_picture(io.BytesIO(chart_images_dict[chart_type]), left, top, width, height)
                charts_added += 1
                
                # Add chart label below
                label_top = top + height + Inches(0.1)
                label_textbox = slide.shapes.add_textbox(left, label_top, width, Inches(0.3))
                label_frame = label_textbox.text_frame
                label_frame.text = chart_type.replace('_', ' ').title()
                label_para = label_frame.paragraphs[0]
                label_para.font.size = Pt(11)
                label_para.font.bold = True
                label_para.font.color.rgb = RGBColor(100, 100, 100)
                label_para.alignment = PP_ALIGN.CENTER
            except Exception as e:
                logger.warning(f"Failed to add {chart_type} chart: {e}")
    
    # If no charts were added, add a message
    if charts_added == 0:
        no_charts_textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
        no_charts_frame = no_charts_textbox.text_frame
        no_charts_frame.text = "No charts available for this report."
        no_charts_para = no_charts_frame.paragraphs[0]
        no_charts_para.font.size = Pt(16)
        no_charts_para.font.color.rgb = RGBColor(150, 150, 150)
        no_charts_para.alignment = PP_ALIGN.CENTER
    
    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes.getvalue()


def create_chart_image(chart_type: str, data: Dict, width: int = 800, height: int = 600) -> bytes:
    """
    Create a chart image from data.
    
    Args:
        chart_type: Type of chart ('timeline', 'product', 'source', 'sentiment')
        data: Chart data dictionary
        width: Image width in pixels
        height: Image height in pixels
    
    Returns:
        bytes: PNG image as bytes
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("matplotlib not available")
    
    plt.style.use('seaborn-v0_8-darkgrid')
    fig, ax = plt.subplots(figsize=(width/100, height/100), dpi=100)
    
    if chart_type == 'timeline':
        # Timeline chart
        dates = data.get('dates', [])
        counts = data.get('counts', [])
        ax.bar(dates, counts, color='#0099ff', alpha=0.7)
        ax.set_xlabel('Date', fontsize=10)
        ax.set_ylabel('Number of Posts', fontsize=10)
        ax.set_title('Posts Timeline', fontsize=12, fontweight='bold')
        plt.xticks(rotation=45, ha='right')
        
    elif chart_type == 'product':
        # Product distribution
        products = data.get('products', [])
        counts = data.get('counts', [])
        colors = plt.cm.Set3(range(len(products)))
        ax.pie(counts, labels=products, autopct='%1.1f%%', colors=colors, startangle=90)
        ax.set_title('Distribution by Product', fontsize=12, fontweight='bold')
        
    elif chart_type == 'source':
        # Source distribution
        sources = data.get('sources', [])
        counts = data.get('counts', [])
        ax.barh(sources, counts, color='#34d399')
        ax.set_xlabel('Number of Posts', fontsize=10)
        ax.set_title('Posts by Source', fontsize=12, fontweight='bold')
        
    elif chart_type == 'sentiment':
        # Sentiment distribution
        sentiments = data.get('sentiments', [])
        counts = data.get('counts', [])
        colors = ['#ef4444', '#f59e0b', '#10b981']  # Red, Orange, Green
        ax.pie(counts, labels=sentiments, autopct='%1.1f%%', colors=colors[:len(sentiments)], startangle=90)
        ax.set_title('Sentiment Distribution', fontsize=12, fontweight='bold')
    
    plt.tight_layout()
    
    # Save to bytes
    img_bytes = io.BytesIO()
    plt.savefig(img_bytes, format='png', dpi=100, bbox_inches='tight')
    plt.close(fig)
    img_bytes.seek(0)
    
    return img_bytes.getvalue()


def prepare_chart_data(posts: List[Dict]) -> Dict:
    """
    Prepare chart data from posts for PowerPoint generation.
    
    Args:
        posts: List of post dictionaries
    
    Returns:
        Dict with chart data for timeline, product, source, and sentiment
    """
    chart_data = {}
    
    if not posts:
        return chart_data
    
    # Sentiment Distribution
    sentiment_counts = Counter()
    for post in posts:
        sentiment = post.get('sentiment_label', 'neutral') or 'neutral'
        sentiment_counts[sentiment] += 1
    
    if sentiment_counts:
        chart_data['sentiment'] = {
            'sentiments': list(sentiment_counts.keys()),
            'counts': [sentiment_counts[s] for s in sentiment_counts.keys()]
        }
    
    # Timeline Chart - Group by date
    timeline_groups = defaultdict(int)
    for post in posts:
        created_at = post.get('created_at', '')
        if created_at:
            try:
                # Parse date
                if 'T' in created_at:
                    date_str = created_at.split('T')[0]
                else:
                    date_str = created_at.split()[0]
                timeline_groups[date_str] += 1
            except:
                continue
    
    if timeline_groups:
        sorted_dates = sorted(timeline_groups.keys())
        chart_data['timeline'] = {
            'dates': sorted_dates,
            'counts': [timeline_groups[d] for d in sorted_dates]
        }
    
    # Source Distribution
    source_counts = Counter()
    for post in posts:
        source = post.get('source', 'Unknown') or 'Unknown'
        source_counts[source] += 1
    
    if source_counts:
        # Sort by count descending, take top 10
        top_sources = source_counts.most_common(10)
        chart_data['source'] = {
            'sources': [s[0] for s in top_sources],
            'counts': [s[1] for s in top_sources]
        }
    
    # Product Distribution - Simple keyword-based detection
    product_patterns = [
        ('VPS', r'\b(vps|virtual\s*private\s*server)\b'),
        ('Dedicated', r'\b(dedicated|dédié|bare\s*metal)\b'),
        ('Public Cloud', r'\b(public\s*cloud|openstack|instance)\b'),
        ('Private Cloud', r'\b(private\s*cloud|vmware)\b'),
        ('Hosting', r'\b(web\s*host|hosting|hébergement)\b'),
        ('Domain', r'\b(domain|domaine|dns)\b'),
        ('Email', r'\b(email|exchange|mail|mx)\b'),
        ('Storage', r'\b(object\s*storage|swift|s3)\b'),
        ('CDN', r'\b(cdn|content\s*delivery)\b'),
        ('Backup', r'\b(backup|veeam|archive)\b'),
    ]
    
    product_counts = Counter()
    for post in posts:
        content = (post.get('content', '') or '').lower()
        detected = False
        for product_name, pattern in product_patterns:
            if re.search(pattern, content, re.IGNORECASE):
                product_counts[product_name] += 1
                detected = True
                break  # Only count first match
        if not detected:
            product_counts['Other'] += 1
    
    if product_counts:
        # Take top 8 products
        top_products = product_counts.most_common(8)
        chart_data['product'] = {
            'products': [p[0] for p in top_products],
            'counts': [p[1] for p in top_products]
        }
    
    return chart_data

